"""python-pptx PowerPoint export renderer."""

from __future__ import annotations

import io
from datetime import datetime, timezone

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ...models.schemas import ParsedPaper
from ..pdf_parser import resolve_figure_path
from .content import SECTION_LABELS, gather_export_context, slugify_title
from .export_formatters import (
    assumptions_bullets,
    cross_qa_entries,
    notes_bullets,
    prepare_sections,
    qa_entries,
    related_bibliography,
    selection_entries,
    summary_sections,
)

_THEMES = {
    "light": {"bg": (255, 255, 255), "text": (24, 24, 27), "muted": (82, 82, 91), "rule": (228, 228, 231)},
    "dark": {"bg": (24, 24, 27), "text": (250, 250, 250), "muted": (161, 161, 170), "rule": (63, 63, 70)},
}


def _rgb(theme: dict, key: str):
    from pptx.dml.color import RGBColor

    return RGBColor(*theme[key])


def _set_slide_bg(slide, theme: dict) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme, "bg")


def _blank_slide(prs: Presentation, theme: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme)
    return slide


def _add_title_slide(prs: Presentation, paper: ParsedPaper, theme: dict) -> None:
    slide = _blank_slide(prs, theme)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = paper.title or "Untitled"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme, "text")
    p.alignment = PP_ALIGN.LEFT
    if paper.authors:
        sub = tf.add_paragraph()
        sub.text = ", ".join(paper.authors)
        sub.font.size = Pt(14)
        sub.font.color.rgb = _rgb(theme, "muted")
    foot = tf.add_paragraph()
    foot.text = f"Analysis export · {datetime.now(timezone.utc).strftime('%B %d, %Y')}"
    foot.font.size = Pt(11)
    foot.font.color.rgb = _rgb(theme, "muted")
    foot.space_before = Pt(18)


def _add_section_title_slide(prs: Presentation, title: str, theme: dict) -> None:
    slide = _blank_slide(prs, theme)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme, "text")


def _add_bullet_slide(
    prs: Presentation,
    heading: str,
    bullets: list[str],
    theme: dict,
    *,
    dense: bool = False,
    subtitle: str | None = None,
) -> None:
    if not bullets:
        return
    chunk = 7 if dense else 5
    for i in range(0, len(bullets), chunk):
        slide = _blank_slide(prs, theme)
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(12), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading if i == 0 else f"{heading} (continued)"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = _rgb(theme, "text")
        if subtitle and i == 0:
            sp = tf.add_paragraph()
            sp.text = subtitle
            sp.font.size = Pt(12)
            sp.font.color.rgb = _rgb(theme, "muted")
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.8), Inches(5.4))
        btf = body.text_frame
        btf.word_wrap = True
        for j, b in enumerate(bullets[i : i + chunk]):
            para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            para.text = b[:900]
            para.font.size = Pt(13 if dense else 15)
            para.font.color.rgb = _rgb(theme, "text")
            para.space_after = Pt(8)


def _add_qa_slide(prs: Presentation, heading: str, pairs: list[tuple[str, str]], theme: dict) -> None:
    for q, a in pairs[:24]:
        slide = _blank_slide(prs, theme)
        box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(12), Inches(6.2))
        tf = box.text_frame
        tf.word_wrap = True
        hp = tf.paragraphs[0]
        hp.text = heading
        hp.font.size = Pt(14)
        hp.font.bold = True
        hp.font.color.rgb = _rgb(theme, "muted")
        qp = tf.add_paragraph()
        qp.text = q[:500]
        qp.font.size = Pt(18)
        qp.font.bold = True
        qp.font.color.rgb = _rgb(theme, "text")
        qp.space_before = Pt(10)
        ap = tf.add_paragraph()
        ap.text = a[:1200] if a else "—"
        ap.font.size = Pt(14)
        ap.font.color.rgb = _rgb(theme, "text")
        ap.space_before = Pt(12)


def _add_selection_slide(
    prs: Presentation,
    heading: str,
    rows: list[tuple[str, str, str]],
    theme: dict,
) -> None:
    for action, sel, body in rows[:20]:
        slide = _blank_slide(prs, theme)
        box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(12), Inches(6.2))
        tf = box.text_frame
        tf.word_wrap = True
        hp = tf.paragraphs[0]
        hp.text = f"{heading} · {action.title()}"
        hp.font.size = Pt(13)
        hp.font.color.rgb = _rgb(theme, "muted")
        qp = tf.add_paragraph()
        qp.text = f"“{sel[:420]}”" if sel else "Selected passage"
        qp.font.size = Pt(15)
        qp.font.italic = True
        qp.font.color.rgb = _rgb(theme, "text")
        qp.space_before = Pt(8)
        bp = tf.add_paragraph()
        bp.text = body[:1400] if body else "—"
        bp.font.size = Pt(14)
        bp.font.color.rgb = _rgb(theme, "text")
        bp.space_before = Pt(14)


def render_pptx(export_row: dict, paper: ParsedPaper, cache: dict) -> tuple[bytes, str, str]:
    options = export_row.get("options") or {}
    pptx_opts = options.get("pptx") or {}
    theme_name = pptx_opts.get("theme", "light")
    dense = pptx_opts.get("dense", False)
    theme = _THEMES.get(theme_name, _THEMES["light"])

    sections = export_row.get("sections") or []
    content = gather_export_context(paper, export_row["user_id"], sections)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_title_slide(prs, paper, theme)

    for key in sections:
        label = SECTION_LABELS.get(key, key)
        _add_section_title_slide(prs, label, theme)

        if key == "summary":
            blocks = summary_sections(content)
            if not blocks:
                _add_bullet_slide(prs, label, ["No summary generated yet."], theme, dense=dense)
            for sub, bullets in blocks:
                _add_bullet_slide(prs, sub, bullets, theme, dense=dense)

        elif key == "prepare":
            blocks = prepare_sections(content)
            if not blocks:
                _add_bullet_slide(prs, label, ["No prepare analysis yet."], theme, dense=dense)
            for sub, bullets in blocks:
                _add_bullet_slide(prs, sub, bullets, theme, dense=dense)

        elif key == "assumptions":
            bullets = assumptions_bullets(content)
            _add_bullet_slide(
                prs, label, bullets or ["No assumptions extracted yet."], theme, dense=dense
            )

        elif key == "qa":
            pairs = qa_entries(content)
            if pairs:
                _add_qa_slide(prs, label, pairs, theme)
            else:
                _add_bullet_slide(prs, label, ["No Q&A yet."], theme, dense=dense)

        elif key == "cross":
            pairs = cross_qa_entries(content)
            if pairs:
                _add_qa_slide(prs, label, pairs, theme)
            else:
                _add_bullet_slide(prs, label, ["No cross-paper Q&A yet."], theme, dense=dense)

        elif key == "notes":
            bullets = notes_bullets(content)
            _add_bullet_slide(
                prs, label, bullets or ["No notes saved yet."], theme, dense=dense
            )

        elif key == "selection":
            rows = selection_entries(content)
            if rows:
                _add_selection_slide(prs, label, rows, theme)
            else:
                _add_bullet_slide(prs, label, ["No selection history yet."], theme, dense=dense)

        elif key == "related":
            bib = related_bibliography(content)
            _add_bullet_slide(
                prs, label, bib or ["No references parsed yet."], theme, dense=True
            )

        elif key == "figures":
            from .export_formatters import figure_slides

            fig_rows = figure_slides(content, paper)
            if not fig_rows:
                _add_bullet_slide(prs, label, ["No figures yet."], theme, dense=dense)
            user_id = export_row["user_id"]
            for cap, analysis, fid in fig_rows:
                slide = _blank_slide(prs, theme)
                path = resolve_figure_path(paper.id, fid, user_id) if fid else None
                if path and path.exists():
                    slide.shapes.add_picture(str(path), Inches(0.6), Inches(0.8), height=Inches(4.8))
                    box = slide.shapes.add_textbox(Inches(6.4), Inches(0.8), Inches(6.2), Inches(5.8))
                else:
                    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.8), Inches(11.8), Inches(5.8))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = cap
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = _rgb(theme, "text")
                if analysis:
                    ap = tf.add_paragraph()
                    ap.text = analysis
                    ap.font.size = Pt(13)
                    ap.font.color.rgb = _rgb(theme, "text")
                    ap.space_before = Pt(10)

    buf = io.BytesIO()
    prs.save(buf)
    slug = slugify_title(paper.title)
    date = datetime.now(timezone.utc).strftime("%Y%m%d")
    filename = f"Know-export-{slug}-{date}.pptx"
    return buf.getvalue(), (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ), filename
