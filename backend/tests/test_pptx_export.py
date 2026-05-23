"""PPTX export smoke tests."""

from __future__ import annotations

import io

from pptx import Presentation

from app.models.schemas import ParsedPaper
from app.services.exports.content import gather_export_context
from app.services.exports.pptx_render import render_pptx


def test_render_pptx_slide_count():
    paper = ParsedPaper(
        id="p1",
        title="Deck Title",
        authors=["Bob"],
        cached_analysis={"summary": {"overview": "Overview text", "key_contributions": ["A", "B"]}},
    )
    export_row = {
        "user_id": "u1",
        "sections": ["summary"],
        "options": {"pptx": {"theme": "light", "dense": False}},
    }
    cache = gather_export_context(paper, "u1", ["summary"])
    data, ctype, filename = render_pptx(export_row, paper, cache)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 2
    assert prs.slides[0].shapes[0].text_frame.text.startswith("Deck Title")
    assert filename.endswith(".pptx")
    assert "presentationml" in ctype


def test_render_pptx_prepare_and_assumptions():
    paper = ParsedPaper(
        id="p2",
        title="Methods Paper",
        authors=["Alice"],
        cached_analysis={
            "pre_reading": {
                "definitions": [{"term": "Loss", "definition": "Objective minimized during training."}],
            },
            "assumptions": {
                "assumptions": [{"type": "explicit", "statement": "Labels are noise-free."}],
            },
        },
    )
    export_row = {
        "user_id": "u1",
        "sections": ["prepare", "assumptions"],
        "options": {"pptx": {"theme": "light", "dense": False}},
    }
    cache = gather_export_context(paper, "u1", export_row["sections"])
    data, _, _ = render_pptx(export_row, paper, cache)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 4
